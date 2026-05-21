#!/usr/bin/env python3
"""Gera o Status Report 1 (sr1_2026-05-22.pptx) reusando os layouts ricos do
template `qualificacao-fmd-exemplo.pptx` (cards G1-G5, 4 fases, dois cards
lado-a-lado, caixa PLANO B).

Aplica o framework narrativo **SCQA** (Situation · Complication · Question ·
Answer) ao longo dos 8 slides:

1. Capa
2. Agenda
3. S + C — Contextualização (problema + gap)         [NOVO slide, duplicado]
4. Q + A — Escopo + Fundamentação (solução + stats)  [layout cards do antigo s3]
5. A detalhada — 5 grupos G1-G5
6. Como temporal — Cronograma 4 marcos M1-M4
7. Then-Now-Next — Status (feito + próximas)
8. Obrigado

Este é um script *custom* para esta apresentação específica.

Uso:
    python3 atividades/apresentacoes/build_sr1_pptx.py
"""
from __future__ import annotations

import sys
from pathlib import Path

# Importar helpers da skill apresentacao-poli
ROOT = Path(__file__).parent.parent.parent
SKILL_SCRIPTS = ROOT / ".claude/skills/apresentacao-poli/scripts"
sys.path.insert(0, str(SKILL_SCRIPTS))

from pptx import Presentation
from pptx_helpers import (  # noqa: E402
    duplicate_slide,
    replace_run_text,
    update_footers,
)

TEMPLATE = ROOT / ".claude/skills/apresentacao-poli/templates/qualificacao-fmd-exemplo.pptx"
OUTPUT = ROOT / "atividades/apresentacoes/sr1_2026-05-22.pptx"


# -------------------------------------------------------------------- helpers


def find_shape_by_text(slide, text: str):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if text in run.text:
                    return shape
    return None


def replace_substring_in_slide(slide, old: str, new: str) -> int:
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    n += run.text.count(old)
                    run.text = run.text.replace(old, new)
    return n


def replace_run_texts_in_order(slide, old: str, new_values: list[str]) -> int:
    idx = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old and idx < len(new_values):
                    run.text = new_values[idx]
                    idx += 1
    return idx


# ------------------------------------------------------------------- por slide


def fix_slide_1_capa(slide) -> None:
    """Capa: equipe sem '(apresentador)', professor, título."""
    replace_run_text(slide, "{{CONTEXTO}}", "OBIONE · STATUS REPORT 1")
    # Título principal (multi-line no original)
    title_shape = find_shape_by_text(slide, "Agentes Conversacionais")
    if title_shape:
        tf = title_shape.text_frame
        first_run = tf.paragraphs[0].runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        font_bold = first_run.font.bold
        try:
            font_color = first_run.font.color.rgb
        except (AttributeError, TypeError):
            font_color = None
        tf.text = "ObiOne · Observatório de Projetos"
        new_run = tf.paragraphs[0].runs[0]
        if font_name:
            new_run.font.name = font_name
        if font_size:
            new_run.font.size = font_size
        if font_bold is not None:
            new_run.font.bold = font_bold
        if font_color:
            new_run.font.color.rgb = font_color
    replace_run_text(
        slide,
        "{{SUBTITULO}}",
        "Backend MVP entregue · 4 semanas antes do previsto",
    )
    replace_run_text(slide, "{{ROTULO_AUTOR}}", "Equipe")
    # Listar os 4 integrantes do grupo
    replace_run_text(
        slide,
        "Raniel Silva",
        "Bruno Rocha · Cynthia Oliveira · Moisés Júnior · Raniel Silva",
    )
    replace_run_text(slide, "Orientador", "Professor")
    replace_run_text(
        slide, "Prof. Dr. Alexandre Maciel", "Prof. Ivaldir Honório de Farias Júnior"
    )
    # Remover linha de programa (já está no subtítulo agora — evita redundância)
    shape = find_shape_by_text(slide, "{{PROGRAMA_LONGO}}")
    if shape:
        shape._element.getparent().remove(shape._element)
    replace_run_text(
        slide,
        "{{INSTITUICAO_LONGA}}",
        "Universidade de Pernambuco · Escola Politécnica de Pernambuco",
    )
    replace_substring_in_slide(slide, "Recife, {{DATA}}", "Recife, Maio 2026")


def fix_slide_2_agenda(slide) -> None:
    """Agenda com 4 itens refletindo o fluxo SCQA + Then-Now-Next."""
    replace_run_text(slide, "{{ROTULO_AGENDA}}", "STATUS REPORT 1")
    replace_run_text(slide, "{{AGENDA_1}}", "Contexto")
    replace_run_text(slide, "{{AGENDA_1_SUB}}", "Por que um observatório vivo")
    replace_run_text(slide, "{{AGENDA_2}}", "O escopo")
    replace_run_text(slide, "{{AGENDA_2_SUB}}", "O que vamos construir")
    replace_run_text(slide, "{{AGENDA_3}}", "Cronograma")
    replace_run_text(slide, "{{AGENDA_3_SUB}}", "9 semanas · 4 marcos")
    replace_run_text(slide, "{{AGENDA_4}}", "Status atual")
    replace_run_text(slide, "{{AGENDA_4_SUB}}", "Onde estamos e o que vem")
    # Remover textbox "Duração estimada"
    shape = find_shape_by_text(slide, "Duração estimada")
    if shape:
        shape._element.getparent().remove(shape._element)


def fix_slide_3_contextualizacao(slide) -> None:
    """[NOVO — slide duplicado] S + C do framework SCQA.

    Card esquerdo: O CONTEXTO (situação atual)
    Card direito: O GAP (complicação — por que mantê-los vivos é difícil)
    """
    replace_run_text(slide, "MOTIVAÇÃO", "CONTEXTO")
    replace_run_text(
        slide, "Por que essa pesquisa?", "Observatório é mais que dashboard"
    )

    # Card esquerda — O CONTEXTO
    replace_run_text(slide, "O CONTEXTO", "OBSERVATÓRIOS VS. FERRAMENTAS")
    replace_run_text(
        slide,
        "Dados industriais subutilizados",
        "Comunidade ativa de múltiplos atores",
    )
    replace_run_text(
        slide,
        "O FMD (Framework de Mineração de Dados; Maciel et al.) define 7 camadas: coleta, governança, AutoML, visualização, deploy, arquitetura e validação.",
        "Um observatório é espaço de conhecimento entre consultoria, clientes e pesquisa, mediado por interação contínua.",
    )
    replace_run_text(
        slide,
        "Hoje, cada camada exige especialista técnico. Engenheiros de chão de fábrica não acessam dados de produção sem ajuda da TI.",
        "Ferramentas de gestão oferecem visualização do status. Faltam o tecido social e a curadoria do conhecimento.",
    )
    replace_run_text(
        slide,
        "Resultado: dados existem, mas decisões não são informadas por eles.",
        "Resultado: observatórios reais são raros porque mantê-los vivos custa caro.",
    )

    # Card direita — O GAP
    replace_run_text(slide, "A LACUNA NA LITERATURA", "POR QUE É CARO?")
    replace_run_text(slide, "Evidência da nossa RSL", "Três fontes de fricção")

    # Stat 1 — Manual
    replace_run_text(slide, "30 / 38", "44 atrib.")
    replace_run_text(
        slide,
        "papers cobrem apenas a camada Coleta",
        "por projeto, curados à mão",
    )
    replace_run_text(
        slide,
        "as 6 outras camadas estão sub-representadas",
        "ler .docx, estruturar, atualizar",
    )

    # Stat 2 — Repetitivo
    replace_run_text(slide, "1 / 38", "Repetitivo")
    replace_run_text(
        slide,
        "paper testou em contexto automotivo",
        "comunicar progresso aos clientes",
    )
    replace_run_text(
        slide,
        "manuais do Ford Fiesta (não dados de produção)",
        "traduzir técnico em narrativa acessível",
    )

    # Stat 3 — Custoso
    replace_run_text(slide, "−53%", "Custoso")
    replace_run_text(
        slide,
        "queda relativa de precisão NL2SQL",
        "manter engajamento da comunidade",
    )
    replace_run_text(
        slide,
        "86,6% no Spider → 41% no banco real Petrobras",
        "responder, contextualizar, sugerir próximos passos",
    )

    # Conclusão rodapé
    replace_run_text(
        slide,
        "Objetivo da tese: framework multi-agente para o FMD, validado em planta automotiva real (Stellantis).",
        "Sem reduzir essa fricção, o observatório fica restrito à teoria. É aí que entra a IA Generativa.",
    )

    # Fontes
    replace_run_text(
        slide,
        "Fonte: Síntese cruzada dos 38 papers da RSL",
        "Fonte: Vieira (2022) · Cap. 5 (MPO)",
    )
    replace_run_text(
        slide,
        "Fonte: Medeiros et al. (2023) · LLM chatbot Ford Fiesta · MDPI Vehicles",
        "Fonte: OPTI-PE como caso prático",
    )
    replace_run_text(
        slide,
        "Fonte: Gao et al. (2024) · Nascimento et al. (2025)",
        "Fonte: Vieira (2022) · Seção 6.4 · Trabalhos Futuros",
    )


def fix_slide_4_escopo_fundamentacao(slide) -> None:
    """Q + A do framework SCQA.

    Card esquerdo: O QUE PROPOMOS (resposta = ObiOne)
    Card direito: ENDEREÇANDO O MPO (validação com stats fundamentais)
    """
    replace_run_text(slide, "MOTIVAÇÃO", "O QUE É O OBIONE")
    replace_run_text(
        slide, "Por que essa pesquisa?", "Observatório-de-portfólio para consultoria de projetos"
    )

    # Card esquerda — manifesto por persona
    replace_run_text(slide, "O CONTEXTO", "PARA QUEM E PARA QUÊ")
    replace_run_text(slide, "Dados industriais subutilizados", "Consultoria · Clientes · Comunidade")
    replace_run_text(
        slide,
        "O FMD (Framework de Mineração de Dados; Maciel et al.) define 7 camadas: coleta, governança, AutoML, visualização, deploy, arquitetura e validação.",
        "Consultoria: observa o portfólio como conhecimento estruturado, não como pasta de Drive.",
    )
    replace_run_text(
        slide,
        "Hoje, cada camada exige especialista técnico. Engenheiros de chão de fábrica não acessam dados de produção sem ajuda da TI.",
        "Cliente: acompanha o próprio projeto em linguagem acessível, sem depender de reunião.",
    )
    replace_run_text(
        slide,
        "Resultado: dados existem, mas decisões não são informadas por eles.",
        "Comunidade vira ativo: comentários, comparativos cross-projeto, IA mantém o tecido vivo.",
    )

    # Card direita — 3 diferenciais (vs.)
    replace_run_text(slide, "A LACUNA NA LITERATURA", "DIFERENCIAL")
    replace_run_text(slide, "Evidência da nossa RSL", "Combinação inédita em 3 vetores")

    replace_run_text(slide, "30 / 38", "vs. PM")
    replace_run_text(
        slide,
        "papers cobrem apenas a camada Coleta",
        "Jira, Trello: gerenciam tarefas e prazos",
    )
    replace_run_text(
        slide,
        "as 6 outras camadas estão sub-representadas",
        "ObiOne observa o projeto como objeto (MPO · 44 atributos)",
    )

    replace_run_text(slide, "1 / 38", "vs. BI")
    replace_run_text(
        slide, "paper testou em contexto automotivo", "Power BI, Looker: dashboards estáticos"
    )
    replace_run_text(
        slide,
        "manuais do Ford Fiesta (não dados de produção)",
        "ObiOne extrai significado de .docx com LLM",
    )

    replace_run_text(slide, "−53%", "vs. acad")
    replace_run_text(
        slide,
        "queda relativa de precisão NL2SQL",
        "OPTI-PE e similares: instrumentos de pesquisa",
    )
    replace_run_text(
        slide,
        "86,6% no Spider → 41% no banco real Petrobras",
        "ObiOne é operacional: consultoria + clientes no dia-a-dia",
    )

    # Conclusão
    replace_run_text(
        slide,
        "Objetivo da tese: framework multi-agente para o FMD, validado em planta automotiva real (Stellantis).",
        "Único a combinar observação MPO + extração com IA + comunidade ativa em uma só plataforma.",
    )

    # Fontes
    replace_run_text(
        slide,
        "Fonte: Síntese cruzada dos 38 papers da RSL",
        "Base: Quadro 37 (Vieira, 2022 · p. 264)",
    )
    replace_run_text(
        slide,
        "Fonte: Medeiros et al. (2023) · LLM chatbot Ford Fiesta · MDPI Vehicles",
        "Referência: OPTI-PE (Vieira, 2022 · Cap. 5)",
    )
    replace_run_text(
        slide,
        "Fonte: Gao et al. (2024) · Nascimento et al. (2025)",
        "Análise: levantamento de ferramentas de mercado",
    )


def fix_slide_5_cinco_grupos(slide) -> None:
    """A detalhada — 5 grupos G1-G5 com textos curtos (≤90 chars).
    Inclui a evidência empírica do smoke Ollama no card Pipeline LLM (G1).
    """
    replace_run_text(slide, "RESULTADOS", "ESCOPO")
    replace_run_text(
        slide,
        "Cinco lacunas identificadas na literatura",
        "5 grupos · 18 requisitos funcionais · smoke real Ollama em 5/5 docs",
    )
    # G1
    replace_run_text(slide, "Governança e Deploy", "Pipeline LLM")
    replace_run_text(
        slide,
        "Nenhum paper dedicado a conformidade LGPD/GDPR ou monitoramento de modelos em produção com interface conversacional",
        "Extração dos 44 atributos via Llama 3.1 8B. Smoke em 5/5 projetos · 46% cobertura média em ~66s/doc.",
    )
    # G2
    replace_run_text(slide, "Integração entre camadas", "Observação")
    replace_run_text(
        slide,
        "Nenhum framework cobre as 7 camadas de forma unificada. O máximo encontrado foi 5/7 camadas em um único trabalho (Keskin et al., 2025)",
        "Dashboard com cobertura do MPO. Cliente vê apenas o seu projeto.",
    )
    # G3
    replace_run_text(slide, "Benchmark vs. realidade", "Comunidade")
    replace_run_text(
        slide,
        "Performance cai drasticamente em bancos reais: 86,6% no Spider → 41% no banco industrial da Petrobras (46 pontos · Gao et al., 2024 · Nascimento et al., 2025)",
        "Autenticação, perfis semi-abertos, comentários e feed in-app.",
    )
    # G4
    replace_run_text(slide, "Pipeline vs. agentes", "IA-Assistente")
    replace_run_text(
        slide,
        "55% dos papers usam pipelines fixos sem capacidade de planejamento, memória ou uso de ferramentas externas",
        "Resumo do projeto (tradutora) + drafts assistidos (redutora de fricção).",
    )
    # G5
    replace_run_text(slide, "Validação automotiva", "Avaliação")
    replace_run_text(
        slide,
        "Apenas 1 de 38 papers testou em contexto automotivo (manuais do Ford Fiesta · Medeiros et al., 2023), escopo muito limitado",
        "Precisão · recall · F1 · concordância humana em 3 projetos + Likert × 2.",
    )


def _fix_5_card_slide(
    slide,
    top_label: str,
    subtitle: str,
    cards: list[tuple[str, str]],
) -> None:
    """Helper genérico para customizar um clone do slide-5-grupos.

    `cards` é uma lista de exatamente 5 tuplas (título, corpo) substituindo,
    em ordem, os 5 cards do template (G1-G5). Usa os textos do template como
    string-de-busca; assume que o clone ainda não foi customizado.
    """
    template_titles = [
        "Governança e Deploy",
        "Integração entre camadas",
        "Benchmark vs. realidade",
        "Pipeline vs. agentes",
        "Validação automotiva",
    ]
    template_bodies = [
        "Nenhum paper dedicado a conformidade LGPD/GDPR ou monitoramento de modelos em produção com interface conversacional",
        "Nenhum framework cobre as 7 camadas de forma unificada. O máximo encontrado foi 5/7 camadas em um único trabalho (Keskin et al., 2025)",
        "Performance cai drasticamente em bancos reais: 86,6% no Spider → 41% no banco industrial da Petrobras (46 pontos · Gao et al., 2024 · Nascimento et al., 2025)",
        "55% dos papers usam pipelines fixos sem capacidade de planejamento, memória ou uso de ferramentas externas",
        "Apenas 1 de 38 papers testou em contexto automotivo (manuais do Ford Fiesta · Medeiros et al., 2023), escopo muito limitado",
    ]
    if len(cards) != 5:
        raise ValueError(f"esperado 5 cards, recebi {len(cards)}")
    replace_run_text(slide, "RESULTADOS", top_label)
    replace_run_text(slide, "Cinco lacunas identificadas na literatura", subtitle)
    # Já validamos len(cards) == 5 acima; zip simples basta. (zip(strict=) só em 3.10+)
    for (new_title, new_body), tpl_title, tpl_body in zip(
        cards, template_titles, template_bodies
    ):
        replace_run_text(slide, tpl_title, new_title)
        replace_run_text(slide, tpl_body, new_body)


def fix_slide_5a_fundacao_pipeline(slide) -> None:
    """5a — Fundação + Pipeline LLM. 6 RFs em 5 cards (último combina RF05+RF06).
    Bodies enxutos: o orador detalha; o slide ancora.
    """
    _fix_5_card_slide(
        slide,
        top_label="REQUISITOS · FUNDAÇÃO + PIPELINE",
        subtitle="6 RFs · sprints 2 e 3 · todos implementados",
        cards=[
            (
                "RF01 · Autenticar",
                "JWT 24h. Admin cria usuários, sem cadastro público.",
            ),
            (
                "RF02 · Perfis",
                "Consultor, cliente, admin. Acesso semi-aberto conforme MPO.",
            ),
            (
                "RF03 · Cadastro",
                "Projeto com nome, domínio e descrição.",
            ),
            (
                "RF04 · Upload .docx",
                "Validação de tipo + checksum + limite 50 MB. Rejeita duplicatas.",
            ),
            (
                "RF05 + RF06 · Extração",
                "44 atributos do Quadro 37 via LLM. Prompt e modelo registrados.",
            ),
        ],
    )


def fix_slide_5b_observacao_comunidade(slide) -> None:
    """5b — Observação + Comunidade + IA. 7 RFs em 5 cards."""
    _fix_5_card_slide(
        slide,
        top_label="REQUISITOS · OBSERVAÇÃO + COMUNIDADE + IA",
        subtitle="7 RFs · sprints 3 e 4 · todos implementados",
        cards=[
            (
                "RF07 + RF09 · Portfólio",
                "Visão do consultor com status e cobertura % do MPO.",
            ),
            (
                "RF08 · Detalhe",
                "Documentos, extrações, cobertura e comentários recentes do projeto.",
            ),
            (
                "RF10 + RF11 · Diálogo",
                "Comentários com 1 nível de threading + feed in-app cronológico.",
            ),
            (
                "RF12 · Resumo Cliente",
                "IA tradutora. Sempre revisado pelo consultor antes de publicar.",
            ),
            (
                "RF13 · Drafts",
                "IA redutora de fricção. Próximos passos e pontos de atenção.",
            ),
        ],
    )


def fix_slide_5c_avaliacao(slide) -> None:
    """5c — Avaliação DSR. 5 RFs em 5 cards, mapeamento direto."""
    _fix_5_card_slide(
        slide,
        top_label="REQUISITOS · AVALIAÇÃO DSR",
        subtitle="Sprint 5 · backend pronto · coleta humana pendente",
        cards=[
            (
                "RF14 · Gabarito",
                "Gabarito humano carregado e validado contra os 44 atributos do MPO.",
            ),
            (
                "RF15 · Comparação",
                "Comparação exata em estruturado + rubrica humana 0/0,5/1 + concordância humana.",
            ),
            (
                "RF16 · Likert Consultoria",
                "4 dimensões: drafts, fricção, qualidade do resumo, manutenibilidade.",
            ),
            (
                "RF17 · Likert Clientes",
                "4 dimensões: clareza, utilidade do espaço, diálogo, sentido de inclusão.",
            ),
            (
                "RF18 · Exportação",
                "JSON (research) + CSV long-format (rubrica). Subsidia o relato.",
            ),
        ],
    )


def fix_slide_5d_rnfs(slide) -> None:
    """5d — RNFs (9 não funcionais agrupados em 5 temas)."""
    _fix_5_card_slide(
        slide,
        top_label="REQUISITOS NÃO-FUNCIONAIS",
        subtitle="9 RNFs por tema · todos endereçados arquiteturalmente",
        cards=[
            (
                "Segurança & LGPD",
                "Autenticação obrigatória (RNF01). NDA + criptografia + logs (RNF02).",
            ),
            (
                "Reprodutibilidade",
                "Prompt e modelo LLM registrados (RNF03). Schema JSON versionado (RNF04).",
            ),
            (
                "Performance",
                "Extração LLM ≤ 90s, medido 46-79s (RNF05). REST ≤ 500ms p95 (RNF06).",
            ),
            (
                "Usabilidade",
                "Linguagem cidadã no Resumo do Cliente. Revisão obrigatória (RNF07).",
            ),
            (
                "Manutenibilidade & Custo",
                "Ports & adapters para LLM e storage (RNF08). Ollama local zero-cost (RNF09).",
            ),
        ],
    )


def fix_slide_6_marcos(slide) -> None:
    """Cronograma 4 marcos sem menções pessoais."""
    replace_run_text(slide, "PROPOSTA DE PESQUISA", "CRONOGRAMA")
    replace_run_text(
        slide,
        "Plano de execução em três fases",
        "9 semanas restantes · 4 marcos · entrega 10/07/2026",
    )

    # M1 PREPARAÇÃO
    replace_run_text(slide, "FASE 1", "M1 ✓")
    replace_run_text(slide, "MVP", "PREPARAÇÃO")
    replace_run_texts_in_order(slide, "6–8 meses", ["22-28 mai", "29 mai - 11 jun"])
    replace_run_text(slide, "Orquestrador", "Atributos-alvo")
    replace_run_text(slide, "+ Coletor (NL2SQL)", "+ Protocolo de avaliação")
    replace_run_text(slide, "+ Visualizador (NL2VIS)", "+ Início dos gabaritos (3 projetos)")
    replace_run_text(
        slide,
        "Técnicas mais maduras na literatura (30 papers).",
        "Concluído. Destravou o Sprint 1.",
    )
    replace_run_text(slide, "Artigo em conferência", "Status Report 1 (hoje)")

    # M2 PIPELINE
    replace_run_text(slide, "FASE 2", "M2 ✓")
    replace_run_text(slide, "Contribuição original", "PIPELINE")
    replace_run_text(slide, "Governança (LGPD)", "Cadastro · upload · autenticação")
    replace_run_text(slide, "+ AutoML", "+ Pipeline LLM nos 5 projetos")
    replace_run_text(
        slide,
        "Ataca os gaps onde não há publicações dedicadas.",
        "Backend antecipado. Smoke Ollama 5/5.",
    )
    replace_run_text(slide, "Artigo em journal", "Extrações dos 5 projetos")

    # M3 DASHBOARD + IA
    replace_run_text(slide, "FASE 3", "M3")
    replace_run_text(slide, "Validação", "DASHBOARD + IA")
    replace_run_text(slide, "8–10 meses", "12-25 jun")
    replace_run_text(slide, "Deploy + Validador", "Cobertura · perfis · comentários")
    replace_run_text(slide, "+ Estudo de caso", "+ Resumo Cliente (IA)")
    replace_run_text(slide, "Stellantis (planta automotiva)", "+ Drafts assistidos (IA)")
    replace_run_text(
        slide,
        "Aproveita parceria institucional do grupo.",
        "Onde a comunidade vive.",
    )
    replace_run_text(slide, "Capítulos da tese", "Status Report 2 (19/06)")

    # M4 AVALIAÇÃO
    replace_run_text(slide, "ESCRITA", "M4")
    replace_run_text(slide, "Tese e defesa", "AVALIAÇÃO")
    replace_run_text(slide, "4–6 meses", "26 jun - 2 jul")
    replace_run_text(slide, "Redação Cap. 3-7", "Precisão · Recall · F1 · Kappa")
    replace_run_text(slide, "Revisão por orientador", "+ Likert × 2 audiências")
    replace_run_text(slide, "Defesa pública", "+ Exportação consolidada")
    replace_run_text(
        slide,
        "Buffer para correções e ajustes pós-banca.",
        "Fecha o ciclo DSR.",
    )
    replace_run_text(slide, "Tese defendida", "Relato + Apresentação")

    # Total (sem menção individual)
    replace_run_text(
        slide,
        "Total: 24–32 meses (deadline: dezembro 2028 = 32 meses)",
        "Escrita do relato: 3-9 jul · Apresentação final: 10/07",
    )

    # Remover caixa Plano B inteiramente (3 textboxes + 2 rectangles na região L>7.5 T>6)
    from pptx.util import Emu
    texts_to_delete = [
        "PLANO B (FASE 3)",
        "Caso a parceria Stellantis não formalize:",
        "Migrar para outra montadora regional ou planta de fornecedor (Tier 1) automotivo via grupo de pesquisa.",
    ]
    for t in texts_to_delete:
        shape = find_shape_by_text(slide, t)
        if shape:
            shape._element.getparent().remove(shape._element)
    # Remover shapes restantes na região da caixa (retângulos decorativos sem texto)
    for shape in list(slide.shapes):
        try:
            left_in = Emu(shape.left).inches if shape.left else 0
            top_in = Emu(shape.top).inches if shape.top else 0
            if left_in > 7.5 and 6.0 < top_in < 7.0:
                shape._element.getparent().remove(shape._element)
        except Exception:
            pass


def fix_slide_7_estamos_vs_proximos(slide) -> None:
    """Then-Now-Next: ESTÁ FEITO (Sprint 1, esta semana) + PRÓXIMAS SEMANAS."""
    replace_run_text(slide, "PRÓXIMOS PASSOS", "ONDE ESTAMOS · O QUE VEM")
    replace_run_text(
        slide,
        "Decisões e posicionamentos",
        "Status atual e próximas semanas",
    )

    # Card esquerda
    replace_run_text(slide, "JÁ ESTOU FAZENDO", "ESTÁ FEITO")
    replace_run_text(slide, "Independente de aprovação", "Sprint 0 + Sprint 1 antecipada")
    replace_run_text(slide, "Capítulo 2 da tese (RSL)", "Requisitos + protocolo")
    replace_run_text(
        slide,
        "Todos os dados estão consolidados em formato estruturado.",
        "18 RFs + 9 RNFs + rubrica híbrida 0/0,5/1 + Kappa.",
    )
    replace_run_text(slide, "Submissão do artigo da RSL", "Backend MVP entregue")
    replace_run_text(
        slide,
        "Target: IEEE Access (preferência por Qualis maior se viável).",
        "18/18 USs · 287 testes · 32 endpoints · 8 migrations · CI rodando.",
    )
    replace_run_text(slide, "Implementação do MVP (Fase 1)", "Smoke real do pipeline LLM")
    replace_run_text(
        slide,
        "Orquestrador + Coletor + Visualizador.",
        "Llama 3.1 8B local em 5/5 projetos · 46% cobertura média.",
    )

    # Card direita
    replace_run_text(slide, "POSIÇÕES QUE DEFENDO", "PRÓXIMAS SEMANAS")
    replace_run_text(slide, "Quero confirmação", "O que vem agora")
    replace_run_text(slide, "Manter as 5 RQs", "Frontend e gabarito (sprints 2-4)")
    replace_run_text(
        slide,
        "O FMD é integrado por design. RQ1 é estrutural; RQ2/RQ3/RQ5 são técnicas; RQ4 é validação unificadora.",
        "Bruno integra os 32 endpoints. Cynthia + Moisés produzem o gabarito manual.",
    )
    replace_run_text(slide, "RQ4 com a Stellantis", "Sprint 5 destravada")
    replace_run_text(
        slide,
        "Aproveitar a relação institucional do grupo de pesquisa para a Fase 3.",
        "Backend compara via /extractions/evaluation. Rubrica humana + concordância.",
    )
    replace_run_text(slide, "RSL fechada com 38 papers", "Status Report 2 (19/06)")
    replace_run_text(
        slide,
        "Saturação temática atingida. Fechar Capítulo 2.",
        "Frontend pronto · IA-Assistente em uso real · Likert lançado.",
    )


def fix_slide_8_obrigado(slide) -> None:
    """Obrigado — sem menções individuais, links pros artefatos."""
    replace_run_text(
        slide,
        "Resultados da RSL  ·  Framework FMD-Agent  ·  Worklog",
        "github.com/raniel90/obione  ·  endpoints documentados  ·  diagramas de arquitetura",
    )
    replace_run_text(
        slide,
        "{{AUTOR}}  ·  {{PROGRAMA}}  ·  {{INSTITUICAO}}",
        "Equipe ObiOne  ·  UPE/POLI",
    )
    replace_run_text(
        slide,
        "Orientador: {{ORIENTADOR}}  ·  Recife, {{DATA}}",
        "Professor: Prof. Ivaldir Honório de Farias Júnior  ·  Recife, Maio 2026",
    )


# ----------------------------------------------------------------------- main


def add_speaker_notes(prs) -> None:
    """Adiciona briefing como nota do apresentador em cada um dos 12 slides."""
    notes = {
        0: """[Tempo: ~30s]

Saudação inicial. Apresentar o grupo e a cadeira:
"Somos Bruno, Cynthia, Moisés e Raniel, do grupo ObiOne, da cadeira Tópicos Avançados em Engenharia de Software, do prof. Ivaldir."

Gancho: "Esse é o nosso Status Report 1. Vamos te mostrar o que é o ObiOne, o que vamos entregar até 10/07, e onde estamos agora — e adianto: estamos bem mais à frente do que o cronograma original previa."
""",
        1: """[Tempo: ~30s]

Apresente o roteiro em 4 blocos sem detalhar:
"Quatro partes: o contexto que motiva o projeto, o que vamos construir + os requisitos em detalhe, onde estamos hoje (incluindo o backend já entregue), e fechamos com o cronograma e os próximos passos."
""",
        2: """[Tempo: ~2min — slide-chave para criar tensão narrativa]

CARD ESQUERDO — "Observatórios vs. ferramentas":
"Observatório é MAIS que dashboard. É um espaço de conhecimento entre múltiplos atores — consultoria, clientes, pesquisa — mediado por interação contínua. Ferramentas de gestão como Jira ou planilhas só oferecem visualização de status. Falta o tecido social e a curadoria do conhecimento."

TRANSIÇÃO para o card direito:
"Se observatórios são tão valiosos, por que são raros na prática? Porque mantê-los vivos é CARO."

CARD DIREITO — três fontes de fricção:
1. Manual: ler .docx, estruturar, atualizar a curadoria de atributos
2. Repetitivo: comunicar progresso ao cliente em linguagem acessível
3. Custoso: manter engajamento da comunidade — responder, contextualizar

FECHAMENTO: "Sem reduzir essa fricção, observatório fica restrito à teoria. É aí que entra a IA Generativa — e o ObiOne."

Bases citadas: Vieira (2022) Cap. 5 (MPO); OPTI-PE como caso prático; Seção 6.4 sobre Trabalhos Futuros.
""",
        3: """[Tempo: ~3min — NÚCLEO ESTRATÉGICO da apresentação]

CARD ESQUERDO — "O que é o ObiOne":
Diga claramente: "O ObiOne é um observatório-de-portfólio para consultoria de projetos."

PARA QUEM E PARA QUÊ — três audiências:
- CONSULTORIA: observa o portfólio como conhecimento estruturado, não como pasta de Drive.
- CLIENTE: acompanha o PRÓPRIO projeto em linguagem acessível, sem depender de reunião.
- COMUNIDADE vira ativo: comentários, comparativos cross-projeto, IA mantém o tecido vivo.

PAUSA — deixe assentar.

CARD DIREITO — "Por que isso é diferente do que já existe? Combinação inédita em 3 vetores":
- vs. PM (Jira, Trello, Asana): gerenciam tarefas e prazos. ObiOne OBSERVA o projeto como objeto de conhecimento, mapeando 44 atributos do MPO.
- vs. BI (Power BI, Looker): mostram métricas. ObiOne extrai SIGNIFICADO de .docx com LLM.
- vs. OBSERVATÓRIOS ACADÊMICOS (OPTI-PE): são instrumentos de pesquisa estáticos. ObiOne é OPERACIONAL — consultoria e clientes no dia-a-dia.

FECHAMENTO: "Único a combinar observação MPO + extração com IA + comunidade ativa em uma só plataforma."

Bases citadas: Quadro 37 do MPO (44 atributos); OPTI-PE (Vieira 2022 Cap. 5); levantamento de ferramentas de mercado.
""",
        4: """[Tempo: ~1.5min]

Descer ao nível da execução:
"Agora o como — escopo organizado em 5 grupos, 18 requisitos funcionais. Importante: o pipeline LLM já está comprovado — rodamos em 5/5 projetos do estudo com 46% de cobertura média em ~66s por documento."

Aponte para cada grupo brevemente (NÃO leia tudo):
- G1 PIPELINE LLM: extrai os 44 atributos via Llama 3.1 8B local. Smoke real feito.
- G2 OBSERVAÇÃO: dashboard de cobertura — cliente vê apenas o seu projeto
- G3 COMUNIDADE: auth, perfis semi-abertos, comentários e feed
- G4 IA-ASSISTENTE: Resumo do Cliente (tradutora) + Drafts (redutora de fricção)
- G5 AVALIAÇÃO: rubrica + Cohen's Kappa nos 3 projetos + Likert × 2 audiências

TRANSIÇÃO: "Vou agora detalhar cada bloco de requisitos — começando pela Fundação e o Pipeline."
""",
        5: """[Tempo: ~1.5min]

CARD por CARD — Fundação + Pipeline:
- RF01 AUTENTICAR: JWT 24h, admin cria contas, sem cadastro público.
- RF02 PERFIS: 3 papéis (consultor, cliente, admin). Acesso semi-aberto = característica MPO.
- RF03 CADASTRO: 6 domínios suportados (legal, health, sports, branding, gastronomy, other).
- RF04 UPLOAD .DOCX: sha256 + 50 MB max, rejeita duplicatas.
- RF05 + RF06 EXTRAÇÃO: 44 atributos do Quadro 37, prompt e modelo registrados em cada extração.

FECHAMENTO: "Os 6 requisitos da Fundação e do Pipeline estão TODOS implementados no backend."
""",
        6: """[Tempo: ~1.5min]

CARD por CARD — Observação + Comunidade + IA:
- RF07 + RF09 PORTFÓLIO: visão do consultor com status derivado e cobertura % do MPO.
- RF08 DETALHE: view consolidado por projeto.
- RF10 + RF11 DIÁLOGO: comentários com 1 nível de threading + feed in-app cronológico.
- RF12 RESUMO DO CLIENTE: IA TRADUTORA — extração técnica vira narrativa acessível.
- RF13 DRAFTS: IA REDUTORA DE FRICÇÃO — propõe próximos passos e pontos de atenção.

ENFATIZE: "Toda saída da IA passa por revisão do consultor antes de virar pública. O cliente nunca vê draft."
""",
        7: """[Tempo: ~1.5min]

CARD por CARD — Avaliação DSR:
- RF14 GABARITO: backend já carrega; produção do gabarito é responsabilidade de Cynthia e Moisés.
- RF15 COMPARAÇÃO: critério híbrido — comparação exata em estruturados + rubrica humana 0/0,5/1 em texto livre + Cohen's Kappa entre avaliadores.
- RF16 LIKERT CONSULTORIA: 4 dimensões — utilidade dos drafts, fricção, qualidade do resumo, manutenibilidade.
- RF17 LIKERT CLIENTES: 4 dimensões — clareza, utilidade do espaço, qualidade do diálogo, sentido de inclusão.
- RF18 EXPORTAÇÃO: JSON bundle (research) + CSV long-format (rubrica).

DESTAQUE: "Todo o backend de avaliação está pronto. Sprint 5 vira coleta humana — não código."
""",
        8: """[Tempo: ~1min]

5 TEMAS por 9 RNFs:
- SEGURANÇA & LGPD (RNF01 + RNF02): auth obrigatório, criptografia em trânsito, logs de acesso.
- REPRODUTIBILIDADE (RNF03 + RNF04): versão de prompt + modelo LLM registrados; schema JSON versionado.
- PERFORMANCE (RNF05 + RNF06): extração ≤ 90s (medimos 46-79s) e REST ≤ 500ms p95.
- USABILIDADE (RNF07): linguagem cidadã no Resumo do Cliente.
- MANUTENIBILIDADE + CUSTO (RNF08 + RNF09): ports & adapters arquitetural + Ollama local zero-cost.

PONTE: "RNFs não viram US, mas todos estão endereçados arquiteturalmente. Vamos ao cronograma."
""",
        9: """[Tempo: ~1.5min]

"São 9 semanas até 10/07, divididas em 4 marcos. Mas atenção: M1 e parte de M2 já estão atingidos antecipadamente."

Percorra rapidamente os marcos:
- M1 PREPARAÇÃO (22-28 mai · esta semana): atributos, protocolo, primeiros gabaritos. ✅ ATINGIDO.
- M2 PIPELINE (29 mai - 11 jun): núcleo técnico. Cadastro + upload + auth + extração LLM nos 5 projetos. ✅ ATINGIDO (3 semanas antecipado).
- M3 DASHBOARD + IA (12-25 jun): a comunidade ganha vida. Backend pronto; falta o frontend.
- M4 AVALIAÇÃO (26 jun - 2 jul): fecha o ciclo DSR. Precisão, Recall, F1, Kappa, Likert × 2. Backend pronto; falta coleta humana.

FECHAMENTO: "Escrita do relato de 3-9/07. Apresentação final em 10/07."
""",
        10: """[Tempo: ~1.5min]

"Onde estamos AGORA e o que vem nas próximas semanas."

CARD ESQUERDO — ESTÁ FEITO (Sprint 0 + Sprint 1 antecipada):
1. REQUISITOS + PROTOCOLO: 18 RFs, 9 RNFs, rubrica híbrida 0/0,5/1 + Kappa.
2. BACKEND MVP ENTREGUE: 18/18 USs implementadas. 287 testes verdes. CI rodando em cada PR.
3. SMOKE REAL DO PIPELINE LLM: Llama 3.1 8B local em 5/5 projetos do estudo, 46% cobertura média.

CARD DIREITO — PRÓXIMAS SEMANAS:
4. FRONTEND + GABARITO (sprints 2-4): Bruno integra os 32 endpoints; Cynthia + Moisés produzem o gabarito de 3 projetos.
5. SPRINT 5 DESTRAVADA: backend compara via /extractions/evaluation; falta a rubrica humana 0/0,5/1.
6. STATUS REPORT 2 (19/06): frontend pronto, IA-Assistente em uso real, Likert lançado.
""",
        11: """[Tempo: ~30s + perguntas]

"Esse foi nosso Status Report 1. Documentação completa no repositório raniel90/obione — incluindo os 32 endpoints com request/response real (api_responses.md) e 5 diagramas Mermaid da arquitetura. Abrimos para perguntas."

Antecipe perguntas prováveis (responda só se vierem):
- "Como vão validar o LLM?" → rubrica híbrida + 2 avaliadores independentes + Kappa. Backend já compara estruturado.
- "Quem são os clientes?" → consultoria parceira; 5 projetos reais de domínios distintos (jurídico, saúde, esporte, branding, fitoterápicos).
- "E se a IA errar?" → IA gera DRAFT, consultor revisa antes de publicar. Cliente nunca vê não-revisado.
- "Por que LLM e não regex/parser?" → heterogeneidade dos .docx; LLM lida com variação semântica. Comprovado em smoke.
- "Por que 46% de cobertura?" → muitos atributos do MPO são genuinamente ausentes no documento. Sprint 5 quantifica via precisão/recall.
- "Por que tão adiantado?" → IA generativa acelerou implementação; o trabalho intelectual (requisitos, arquitetura, protocolo) virou base sólida pra automação.
""",
    }
    for idx, text in notes.items():
        prs.slides[idx].notes_slide.notes_text_frame.text = text.strip()
    print(f"  {len(notes)} speaker notes adicionados")


def main() -> int:
    if not TEMPLATE.exists():
        print(f"❌ Template não encontrado: {TEMPLATE}", file=sys.stderr)
        return 2

    print(f"Carregando template: {TEMPLATE.name}")
    prs = Presentation(TEMPLATE)

    print(f"Duplicando slide 3 (cards lado-a-lado) para criar S+C (novo slide 3)...")
    duplicate_slide(prs, src_idx=2, dst_idx=3)
    print(f"  Após duplicação SCQA: {len(prs.slides)} slides")

    print(f"Duplicando slide 5 (5-cards) 4× para slides 5a/5b/5c/5d...")
    # Cada duplicação insere o clone na posição indicada e empurra os
    # subsequentes para a direita. Por isso usamos dst_idx crescente (5,6,7,8)
    # com src_idx fixo em 4 — o slide original permanece em 4.
    duplicate_slide(prs, src_idx=4, dst_idx=5)
    duplicate_slide(prs, src_idx=4, dst_idx=6)
    duplicate_slide(prs, src_idx=4, dst_idx=7)
    duplicate_slide(prs, src_idx=4, dst_idx=8)
    print(f"  Total de slides agora: {len(prs.slides)}")

    print("Slide 1: capa...")
    fix_slide_1_capa(prs.slides[0])

    print("Slide 2: agenda (4 itens novos)...")
    fix_slide_2_agenda(prs.slides[1])

    print("Slide 3: S + C — Contextualização (problema + gap)...")
    fix_slide_3_contextualizacao(prs.slides[2])

    print("Slide 4: Q + A — Escopo + Fundamentação MPO...")
    fix_slide_4_escopo_fundamentacao(prs.slides[3])

    print("Slide 5: 5 grupos (overview)...")
    fix_slide_5_cinco_grupos(prs.slides[4])

    print("Slide 5a: Requisitos · Fundação + Pipeline...")
    fix_slide_5a_fundacao_pipeline(prs.slides[5])

    print("Slide 5b: Requisitos · Observação + Comunidade + IA...")
    fix_slide_5b_observacao_comunidade(prs.slides[6])

    print("Slide 5c: Requisitos · Avaliação DSR...")
    fix_slide_5c_avaliacao(prs.slides[7])

    print("Slide 5d: Requisitos Não-Funcionais...")
    fix_slide_5d_rnfs(prs.slides[8])

    print("Slide 6: Cronograma (M1-M4 + PLANO B)...")
    fix_slide_6_marcos(prs.slides[9])

    print("Slide 7: Status (Então-Agora-Próximo)...")
    fix_slide_7_estamos_vs_proximos(prs.slides[10])

    print("Slide 8: obrigado...")
    fix_slide_8_obrigado(prs.slides[11])

    # Substituir {{FOOTER}} textual em todos os slides
    print("\nSubstituindo {{FOOTER}} global...")
    footer_text = "Equipe ObiOne · UPE/POLI · Maio 2026"
    footer_count = 0
    for slide in prs.slides:
        footer_count += replace_run_text(slide, "{{FOOTER}}", footer_text)
    print(f"  {footer_count} footers textuais atualizados")

    print("\nAtualizando numeração de página (total=12)...")
    n = update_footers(prs, total=12)
    print(f"  {n} números de página atualizados")

    # NOTA: speaker notes desativados temporariamente.
    # Com 4 clones do slide-5-grupos, a infraestrutura de notesMaster do
    # python-pptx falha em criar o part corretamente e o Keynote rejeita o
    # arquivo. O conteúdo dos briefings vive em add_speaker_notes() pra
    # uso manual; até resolver o root cause, deixamos os notes vazios.
    # print("\nAdicionando speaker notes (briefing por slide)...")
    # add_speaker_notes(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"\n✓ Salvo: {OUTPUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
