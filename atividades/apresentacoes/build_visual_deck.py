#!/usr/bin/env python3
"""Reconstrói o miolo da apresentação final com layouts visuais (diagramas,
prints reais, cards de número, camada de IA e experimento detalhado),
reaproveitando capa/agenda/obrigado e a identidade UPE/POLI."""
import sys
from pathlib import Path

BASE = Path(__file__).resolve().parents[2] / ".claude/skills/apresentacao-poli"
sys.path.insert(0, str(BASE / "scripts"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx_helpers import update_footers, delete_slide_by_index

HERE = Path(__file__).resolve().parent
PRINTS = HERE / "prints"

NAVY = RGBColor(0x1F, 0x2A, 0x44)
RED = RGBColor(0xE0, 0x26, 0x1E)
GRAY = RGBColor(0x5B, 0x63, 0x72)
MUTED = RGBColor(0x8A, 0x90, 0x9C)
LIGHT = RGBColor(0xF2, 0xF3, 0xF5)
CARD = RGBColor(0xF6, 0xF7, 0xF9)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDD, 0xE1, 0xE7)
FONT = "Arial"


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        for (text, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = bold
            f.italic = italic
            f.color.rgb = color
    return tb


def eyebrow(slide, text):
    txt(slide, 0.55, 0.5, 12.2, 0.3, [[(text.upper(), 12, RED, True, False)]])


def title(slide, text, size=30):
    txt(slide, 0.55, 0.88, 12.2, 1.0, [[(text, size, NAVY, True, False)]])


def subtitle(slide, text):
    txt(slide, 0.55, 1.9, 12.2, 0.4, [[(text, 15.5, GRAY, False, False)]])


def rect(slide, l, t, w, h, fill, line=None, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


def set_text(shp, paras, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.14))
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(size)
            f.bold = bold
            f.color.rgb = color


def arrow(slide, l, t, w, h, color=RED):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def picture(slide, path, l, t, h=None, w=None):
    kw = {}
    if h is not None:
        kw["height"] = Inches(h)
    if w is not None:
        kw["width"] = Inches(w)
    pic = slide.shapes.add_picture(str(path), Inches(l), Inches(t), **kw)
    pic.line.color.rgb = BORDER
    pic.line.width = Pt(1.25)
    return pic


def stat_card(slide, l, t, w, h, big, label, accent):
    rect(slide, l, t, w, h, CARD, line=BORDER)
    rect(slide, l, t, w, 0.12, accent, rounded=False)
    inner = txt(slide, l, t + 0.28, w, h - 0.4,
                [[(big, 42, NAVY, True, False)], [(label, 13, GRAY, False, False)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return inner


def panel(slide, l, t, w, h, header, rows):
    rect(slide, l, t, w, h, CARD, line=BORDER)
    tb = slide.shapes.add_textbox(Inches(l + 0.25), Inches(t + 0.22), Inches(w - 0.5), Inches(h - 0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = header.upper()
    r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = RED
    p.space_after = Pt(10)
    for head, desc in rows:
        ph = tf.add_paragraph(); ph.space_before = Pt(6); ph.space_after = Pt(1)
        rh = ph.add_run(); rh.text = head
        rh.font.name = FONT; rh.font.size = Pt(14.5); rh.font.bold = True; rh.font.color.rgb = NAVY
        pd = tf.add_paragraph(); pd.space_after = Pt(4)
        rd = pd.add_run(); rd.text = desc
        rd.font.name = FONT; rd.font.size = Pt(12.5); rd.font.color.rgb = GRAY


def clean_content_slide(slide):
    for sh in list(slide.shapes):
        is_pic = sh.shape_type == 13
        top_in = sh.top / 914400 if sh.top is not None else 0
        if (is_pic and top_in < 0.6) or top_in > 6.8:
            continue
        sh._element.getparent().remove(sh._element)


# ─────────────────────────────────────────────────────────────────────────────
def slide_problema(s):
    eyebrow(s, "O ponto de partida")
    title(s, "O conhecimento dos projetos se perde")
    subtitle(s, "Consultorias vivem de conhecimento, e ele evapora entre um projeto e o outro.")
    cards = [
        ("Fica na cabeça das pessoas", "O que um projeto ensinou raramente chega ao próximo"),
        ("Concentrado em poucos", "A conexão entre projetos depende dos consultores mais experientes"),
        ("O cliente participa pouco", "E abrir tudo comprometeria a confidencialidade entre clientes"),
    ]
    w, gap, t, h = 3.85, 0.28, 3.0, 2.5
    x = 0.55
    for head, sub in cards:
        c = rect(s, x, t, w, h, CARD, line=BORDER)
        set_text(c, [[(head, 16, NAVY, True)], [(sub, 12.5, GRAY, False)]])
        x += w + gap


def slide_solucao(s):
    eyebrow(s, "O que construímos")
    title(s, "O ciclo do conhecimento no observatório")
    subtitle(s, "Observações viram discussões na comunidade e se consolidam em conhecimento, ancorado no MPO (Vieira, 2022).")
    steps = ["Observação", "Discussão", "Conhecimento"]
    w, t, h = 3.3, 3.1, 1.25
    xs = [0.9, 5.0, 9.1]
    for i, (x, st) in enumerate(zip(xs, steps)):
        b = rect(s, x, t, w, h, NAVY)
        set_text(b, [[(st, 18, WHITE, True)]])
        if i < 2:
            arrow(s, x + w + 0.02, t + 0.42, 0.62, 0.42)
    band = rect(s, 0.9, 5.05, 11.5, 0.95, LIGHT, line=BORDER)
    set_text(band, [[("O conhecimento consolidado fica disponível para os próximos projetos do domínio", 15, NAVY, True)]])


def slide_ia_layer(s):
    eyebrow(s, "A camada de IA")
    title(s, "Assistiva, ancorada e auditável")
    subtitle(s, "A IA acelera cada ponta do ciclo, sempre com o consultor no comando.")
    panel(s, 0.55, 2.5, 5.9, 4.05, "Papéis da IA no ciclo", [
        ("Configuradora", "No cadastro: propõe o domínio e o que acompanhar"),
        ("Observadora", "Sugere observações à luz do MPO"),
        ("Sintetizadora", "Consolida o aprendizado da conversa"),
        ("Conectora", "Sintetiza padrões entre projetos do domínio"),
    ])
    panel(s, 6.9, 2.5, 5.9, 4.05, "Como confiamos", [
        ("A IA propõe, o humano decide", "Nunca escreve direto: o consultor revisa e publica"),
        ("Ancorada no MPO, com saída estruturada", "A resposta segue um formato definido e usa apenas atributos do catálogo"),
        ("Proveniência registrada", "Cada sugestão guarda provedor, modelo e momento"),
    ])



def slide_pivo(s):
    eyebrow(s, "A concepção")
    title(s, "Da proposta inicial ao observatório")
    subtitle(s, "A crítica dos orientadores levou a um reescopo: refinar o propósito do artefato.")
    b1 = rect(s, 0.9, 3.0, 4.6, 1.6, LIGHT, line=BORDER)
    set_text(b1, [[("Proposta inicial", 13, RED, True)], [("Extrator de atributos do MPO a partir de documentos", 14, NAVY, True)]])
    arrow(s, 5.7, 3.55, 0.9, 0.5)
    b2 = rect(s, 6.85, 3.0, 5.55, 1.6, NAVY)
    set_text(b2, [[("Proposta refinada", 13, WHITE, False)], [("Observatório de projetos com participação do cliente na comunidade", 14.5, WHITE, True)]])
    txt(s, 0.9, 5.1, 11.5, 0.9,
        [[("O foco deslocou-se da extração de documentos para o ciclo de conhecimento, e a pergunta de pesquisa passou a ser: a IA generativa viabiliza um observatório de projetos?", 14.5, GRAY, False, False)]])


def slide_desafios(s):
    eyebrow(s, "Principais desafios")
    title(s, "O que enfrentamos e como resolvemos")
    panel(s, 0.55, 2.0, 5.9, 4.75, "Escopo, modelo e IA", [
        ("Escopo amplo e pivô no meio do projeto", "Reescopo formal com os orientadores e disciplina de cronograma"),
        ("Entender o MPO a fundo e traduzi-lo sem jargão", "Catálogo canônico dos 44 atributos e linguagem simples na interface"),
        ("Confiar na IA sem perder o controle", "Saída estruturada, grounding no MPO, validação em código e revisão humana"),
    ])
    panel(s, 6.9, 2.0, 5.9, 4.75, "Validação, adoção e equipe", [
        ("Validar com usuários reais a distância", "Acesso remoto por túnel e dados de demonstração replicáveis"),
        ("Clareza do primeiro acesso (3,8 nas 2 rodadas)", "Onboarding incluído; o próximo passo é navegação guiada"),
        ("Coordenar 4 frentes em equipe", "Comunicação constante entre extração, front, avaliação e escrita"),
    ])


def slide_shot(s, eb, ttl, img, caption_lines):
    eyebrow(s, eb)
    title(s, ttl)
    picture(s, img, 5.35, 2.15, h=4.75)
    paras = [[("›  ", 15, RED, True, False), (line, 15, NAVY, False, False)] for line in caption_lines]
    txt(s, 0.55, 2.7, 4.5, 4.0, paras, anchor=MSO_ANCHOR.TOP)


def slide_shot_wide(s, eb, ttl, img, caption):
    eyebrow(s, eb)
    title(s, ttl)
    subtitle(s, caption)
    w = 10.4
    picture(s, img, (13.33 - w) / 2, 2.6, w=w)


def slide_experimento(s):
    eyebrow(s, "Como avaliamos")
    title(s, "Método: Design Science Research")
    steps = ["Construir\no artefato", "Demonstrar\nem uso real", "Avaliar\na percepção"]
    w, t, h = 3.2, 2.35, 1.15
    xs = [0.9, 5.06, 9.22]
    for i, (x, st) in enumerate(zip(xs, steps)):
        b = rect(s, x, t, w, h, NAVY)
        set_text(b, [[(seg, 15, WHITE, True)] for seg in st.split("\n")])
        if i < 2:
            arrow(s, x + w + 0.02, t + 0.4, 0.55, 0.36, color=MUTED)
    c1 = rect(s, 0.9, 4.0, 5.55, 1.55, CARD, line=BORDER)
    set_text(c1, [[("Instrumento", 14.5, RED, True)],
                  [("12 afirmações em escala Likert (1 a 5) e 3 perguntas abertas, após demonstração guiada", 13, NAVY, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    c2 = rect(s, 6.85, 4.0, 5.55, 1.55, CARD, line=BORDER)
    set_text(c2, [[("Amostra", 14.5, RED, True)],
                  [("2 rodadas de 4 consultores (N=8), uma consultoria real; 2ª rodada após onboarding", 13, NAVY, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 0.9, 5.75, 11.5, 0.5,
        [[("Perspectiva do consultor. Reportado como casos, sem inferência estatística.", 13, MUTED, False, True)]])


def slide_resultados(s):
    eyebrow(s, "Resultados")
    title(s, "Valor percebido, com um ponto de atenção")
    cards = [
        ("4,3 / 5", "média acumulada (8 consultores)", GREEN),
        ("N = 8", "2 rodadas de 4 · uma consultoria", NAVY),
        ("83%", "das respostas foram nota 4 ou 5", NAVY),
    ]
    w, gap, t, h = 3.85, 0.28, 2.35, 2.15
    x = 0.55
    for big, label, accent in cards:
        stat_card(s, x, t, w, h, big, label, accent)
        x += w + gap
    txt(s, 0.55, 4.75, 12.3, 0.7,
        [[("1ª rodada 4,48; a 2ª foi mais crítica (4,1). Aprendizados e governança no topo "
           "(5,0 na 1ª rodada); a clareza permaneceu em 3,8 e segue como o principal ponto de atenção.", 15, NAVY, False, False)]])
    txt(s, 0.55, 5.85, 12.3, 0.6,
        [[("Nas respostas abertas o valor apareceu no aprendizado consolidado com IA; "
           "as críticas convergem para clareza e navegação.", 13.5, GRAY, False, True)]])


def slide_valor(s):
    eyebrow(s, "Onde chegamos")
    title(s, "O que o ObiOne entrega hoje")
    items = ["Isolamento entre clientes", "Reaproveitamento", "Engajamento do cliente", "IA assistiva"]
    w, gap, t, h = 2.85, 0.22, 3.05, 1.5
    x = 0.55
    for it in items:
        c = rect(s, x, t, w, h, CARD, line=BORDER)
        set_text(c, [[("✓", 20, GREEN, True)], [(it, 14, NAVY, True)]])
        x += w + gap
    txt(s, 0.55, 5.2, 12.2, 0.9,
        [[("Ciclo completo, replicável e demonstrável ao vivo. ", 16, NAVY, True, False),
          ("Próximos passos: navegação guiada e avaliação da fidelidade da extração do MPO.", 16, GRAY, False, False)]])


BUILDERS = [
    slide_problema,
    slide_pivo,
    slide_experimento,
    slide_solucao,
    slide_ia_layer,
    lambda s: slide_shot(s, "O diferencial · acesso semi-aberto", "O cliente participa, sem ver outros clientes",
                         PRINTS / "cliente_novidades.png",
                         ["Vê e acessa apenas o próprio projeto", "Contribui nas conversas da comunidade", "As novidades do projeto o convidam a voltar"]),
    lambda s: slide_shot_wide(s, "Reaproveitamento", "O conhecimento atravessa projetos",
                              PRINTS / "aprendizados_dominio_bloco.png",
                              "Aprendizados consolidados em um projeto aparecem nos demais projetos do domínio; a Conectora sintetiza padrões sob demanda."),
    slide_resultados,
    slide_desafios,
    slide_valor,
]


def sub_runs(prs, replacements):
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in replacements:
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def main():
    src = HERE / "ObiOne_Apresentacao_Final.pptx"
    prs = Presentation(str(src))
    n = len(BUILDERS)
    # garante slots de conteúdo suficientes (capa+agenda = 0-1; último = Obrigado)
    from pptx_helpers import duplicate_slide
    while len(prs.slides._sldIdLst) - 3 < n:
        duplicate_slide(prs, 2, len(prs.slides._sldIdLst) - 1)
    obrigado_idx = len(prs.slides._sldIdLst) - 1
    for i, build in enumerate(BUILDERS):
        s = prs.slides[2 + i]
        clean_content_slide(s)
        build(s)
    # remove slides de conteúdo em excesso (entre o último construído e o Obrigado)
    for idx in range(obrigado_idx - 1, 2 + n - 1, -1):
        delete_slide_by_index(prs, idx)
    sub_runs(prs, [
        ("Doutorado PPGEC · Tópicos Avançados em Engenharia de Software",
         "Tópicos Avançados em Engenharia de Software"),
        ("Doutorado PPGEC", "TAES"),
        ("Discussão e perguntas", "Veremos ao vivo"),
        # agenda alinhada à trajetória (apresentação de projeto, não pitch)
        ("O problema e a solução", "De onde partimos"),
        ("Por que um observatório", "Contexto, concepção e pivô"),
        ("A jornada de valor", "O que construímos"),
        ("O ciclo funcionando", "O ciclo, a IA e a comunidade"),
        ("Diferencial e IA", "Como avaliamos"),
        ("Cliente dentro, IA copiloto", "Método e resultados"),
        ("Validação e valor", "Desafios e demo"),
        ("Evidência e entrega", "O que enfrentamos e o observatório ao vivo"),
        # valores já gravados no pptx committado (idempotência)
        ("A camada de IA", "O que construímos"),
        ("IA assistiva, com revisão humana", "O ciclo, a IA e a comunidade"),
        ("Diferencial e reaproveitamento", "Como avaliamos"),
        ("Cliente com muralha e reuso", "Método e resultados"),
        ("Experimento e valor", "Desafios e demo"),
        ("Método, resultados e entrega", "O que enfrentamos e o observatório ao vivo"),
    ])
    update_footers(prs)
    prs.save(str(src))
    print(f"✓ {src} · {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
