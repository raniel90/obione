#!/usr/bin/env python3
"""QC validator: roda checagens contra os anti-patterns documentados.

Uso:
    python3 qc_check.py <caminho.pptx>

Saída:
    - Lista de violações (slide, tipo, contexto)
    - Exit code 0 se zero issues; 1 se há violações

Validações implementadas:
    1. IDs internos no formato 'autor_AAAA' ou 'autor_palavra_AAAA'
    2. Paths de arquivo (.md, .yaml, rsl_docs/, similar)
    3. Fontes fora do padrão (qualquer != Arial — exceto onde permitido)
    4. Em-dashes (—) — travessões substituem por · ou :
    5. Jargão IA suspeito ('3-skill', 'screening automático', 'Cohen Kappa')
    6. Placeholders não preenchidos ('{{X}}')
"""
from __future__ import annotations

import re
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation

# Patterns de anti-patterns
PAT_PAPER_ID = re.compile(r"\b[a-z][a-z_]*_(19|20)\d{2}\b")
PAT_FILE_PATH = re.compile(r"\.(md|yaml|yml|json)\b|/[a-zA-Z_]+/[a-zA-Z_]+\.[a-z]+")
PAT_EM_DASH = re.compile(r"—")
PAT_AI_JARGON = re.compile(
    r"\b(3.skill|pipeline 3|screening autom[áa]tico|Cohen[' ]?s? [Kk]appa|"
    r"\bκ\s*=|assist[êe]ncia de LLM|NEEDS_HUMAN_REVIEW)\b"
)
PAT_PLACEHOLDER = re.compile(r"\{\{[A-Z_]+\}\}")

# Fontes permitidas além de Arial
FONTS_ALLOWED = {"Arial", None}  # None = herda do master/layout


def scan(pptx_path: Path) -> dict[str, list[tuple[int, str]]]:
    """Retorna {check_name: [(slide_num, contexto), ...]}."""
    prs = Presentation(pptx_path)
    issues: dict[str, list[tuple[int, str]]] = defaultdict(list)

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    txt = run.text
                    fname = run.font.name

                    if PAT_PAPER_ID.search(txt):
                        issues["paper_id"].append((i, _ctx(txt)))
                    if PAT_FILE_PATH.search(txt):
                        issues["file_path"].append((i, _ctx(txt)))
                    if PAT_EM_DASH.search(txt):
                        issues["em_dash"].append((i, _ctx(txt)))
                    if PAT_AI_JARGON.search(txt):
                        issues["ai_jargon"].append((i, _ctx(txt)))
                    if PAT_PLACEHOLDER.search(txt):
                        issues["placeholder"].append((i, _ctx(txt)))
                    if fname not in FONTS_ALLOWED:
                        issues["font"].append((i, f"font={fname!r} text={_ctx(txt)}"))

    return dict(issues)


def _ctx(text: str, n: int = 80) -> str:
    return text.strip()[:n]


CHECK_LABELS = {
    "paper_id": "IDs internos (use 'Autor et al. (Ano)')",
    "file_path": "Refs a arquivos (.md/.yaml/paths) — use nome humano",
    "em_dash": "Travessão (—) — use · ou :",
    "ai_jargon": "Jargão IA não documentado (3-skill, screening automático, κ, etc.)",
    "placeholder": "Placeholder não preenchido ({{X}})",
    "font": "Fonte fora do padrão (esperado: Arial)",
}


def report(issues: dict[str, list[tuple[int, str]]]) -> int:
    """Imprime relatório e retorna exit code (0=ok, 1=falha)."""
    total = sum(len(v) for v in issues.values())
    if total == 0:
        print("✅ PASS — zero anti-patterns detectados.")
        return 0

    print(f"❌ FAIL — {total} anti-pattern(s) detectado(s):\n")
    for check, occurrences in issues.items():
        label = CHECK_LABELS.get(check, check)
        print(f"  [{check}] {label}: {len(occurrences)} ocorrência(s)")
        for slide_num, ctx in occurrences[:5]:
            print(f"     S{slide_num}: {ctx}")
        if len(occurrences) > 5:
            print(f"     ... (+{len(occurrences) - 5} mais)")
        print()
    return 1


def main() -> int:
    if len(sys.argv) != 2:
        print("uso: qc_check.py <caminho.pptx>", file=sys.stderr)
        return 2
    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"arquivo não encontrado: {pptx_path}", file=sys.stderr)
        return 2
    issues = scan(pptx_path)
    return report(issues)


if __name__ == "__main__":
    sys.exit(main())
