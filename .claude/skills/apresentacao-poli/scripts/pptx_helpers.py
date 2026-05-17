"""Helpers reutilizáveis para edição de .pptx no padrão UPE/POLI.

Funções extraídas e generalizadas a partir de scripts/restructure_presentation.py
do projeto phd-fmd-agents-rsl. Nenhuma lógica é específica de RSL/FMD aqui.

Uso típico:
    from pptx import Presentation
    from pptx_helpers import (
        delete_slide_by_index, reorder_slides,
        replace_run_text, delete_shape_by_run_text,
        update_footers, highlight_runs,
    )

    prs = Presentation("template.pptx")
    n = replace_run_text(prs.slides[0], "{{TITULO}}", "Status TAES")
    update_footers(prs, total=len(prs.slides))
    prs.save("output.pptx")
"""
from __future__ import annotations

import re
from typing import Callable, Iterable

from pptx import Presentation as _PresOpen  # noqa: F401 — re-export hint
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


def duplicate_slide(prs, src_idx: int, dst_idx: int | None = None):
    """Duplica o slide na posição `src_idx`. Insere o novo slide em `dst_idx`
    (zero-based) ou no final se `dst_idx is None`.

    Retorna o novo slide. Mantém todos os shapes (incluindo imagens) e a
    relação com o slide layout original.

    Implementação: python-pptx não expõe API de duplicação;
    1. Copia o XML do `spTree` (preserva todas as shapes textuais e decorativas)
    2. Para Pictures, re-adiciona via `add_picture` para garantir que as
       relações `r:embed` apontem para image parts do slide novo (senão a
       imagem aparece quebrada / não renderiza)
    3. Move o `sldId` na `sldIdLst` para `dst_idx` se necessário
    """
    from copy import deepcopy
    from io import BytesIO

    src_slide = prs.slides[src_idx]

    # Coletar pictures do src ANTES de duplicar — vamos re-adicionar com bytes
    pics_to_restore = []
    for shape in src_slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            try:
                pics_to_restore.append({
                    "blob": shape.image.blob,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                })
            except Exception:
                # Picture sem blob acessível (ex: linked, não embedded); skip
                pass

    # Adicionar novo slide com o mesmo layout do src
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    # Remover placeholders padrão
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Copiar todos os shapes do src_slide para o new_slide (inclusive pictures quebradas)
    for shape in src_slide.shapes:
        el = shape._element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.append(new_el)

    # Remover as Pictures duplicadas (broken — rId aponta pra outro slide)
    for shape in list(new_slide.shapes):
        if shape.shape_type == 13:
            shape._element.getparent().remove(shape._element)

    # Re-adicionar Pictures usando os bytes e posições originais
    for pic in pics_to_restore:
        new_slide.shapes.add_picture(
            BytesIO(pic["blob"]),
            pic["left"], pic["top"],
            width=pic["width"], height=pic["height"],
        )

    # Reordenar se dst_idx for especificado e diferente do final
    if dst_idx is not None and dst_idx != len(prs.slides) - 1:
        sld_id_lst = prs.slides._sldIdLst
        sld_ids = list(sld_id_lst)
        new_sld_id = sld_ids[-1]  # último é o que acabamos de adicionar
        sld_id_lst.remove(new_sld_id)
        sld_id_lst.insert(dst_idx, new_sld_id)

    return new_slide


def delete_slide_by_index(prs, idx: int) -> None:
    """Remove o slide na posição zero-based.

    Mexe diretamente no sldIdLst e dropa o rel — python-pptx não expõe API.
    """
    slide_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    sld_ids = list(slide_id_lst)
    sld_id = sld_ids[idx]
    rId = sld_id.attrib[qn("r:id")]
    prs.part.drop_rel(rId)
    slide_id_lst.remove(sld_id)


def reorder_slides(prs, new_order: list[int]) -> None:
    """Reordena slides por índices zero-based atuais.

    new_order deve ter o mesmo tamanho do total de slides.
    """
    slide_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    sld_ids = list(slide_id_lst)
    if len(new_order) != len(sld_ids):
        raise ValueError(
            f"new_order length {len(new_order)} != slide count {len(sld_ids)}"
        )
    for sid in sld_ids:
        slide_id_lst.remove(sid)
    for src_idx in new_order:
        slide_id_lst.append(sld_ids[src_idx])


def replace_run_text(slide, old: str, new: str) -> int:
    """Substitui o texto de runs cujo .text == old (match exato).

    Retorna o número de runs alterados. Preserva formatação (font, cor, bold).
    """
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old:
                    run.text = new
                    n += 1
    return n


def delete_shape_by_run_text(slide, target_text: str) -> int:
    """Deleta todo shape que tenha algum run.text == target_text. Retorna count."""
    to_delete = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == target_text:
                    to_delete.append(shape)
                    break
            else:
                continue
            break
    for shape in to_delete:
        sp = shape._element
        sp.getparent().remove(sp)
    return len(to_delete)


def update_footers(
    prs,
    total: int | None = None,
    pattern: re.Pattern | str = r"^\s*\d+\s*/\s*\d+\s*$",
    new_format: str = "{i} / {total}",
    skip_slides: Iterable[int] = (),
) -> int:
    """Substitui footers no formato 'N / M' pelo padrão atualizado.

    - total: denominador final. Se None, usa len(prs.slides).
    - pattern: regex que identifica o footer. Default casa qualquer 'N / M'.
    - new_format: template com placeholders {i} (posição 1-based) e {total}.
    - skip_slides: índices 1-based de slides sem footer (ex.: capa).

    Retorna número de footers atualizados.
    """
    if isinstance(pattern, str):
        pattern = re.compile(pattern)
    if total is None:
        total = len(prs.slides)
    skip = set(skip_slides)
    n = 0
    for i, slide in enumerate(prs.slides, start=1):
        if i in skip:
            continue
        new_footer = new_format.format(i=i, total=total)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if pattern.match(run.text):
                        run.text = new_footer
                        n += 1
    return n


def highlight_runs(
    slide,
    predicate: Callable[[str], bool],
    bold: bool = True,
    color: RGBColor | None = None,
) -> int:
    """Aplica destaque (bold + cor) em runs onde predicate(run.text) é True.

    Default color é UPE_RED (#E0261E) se nenhuma cor for passada.
    """
    if color is None:
        color = UPE_RED
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if predicate(run.text):
                    if bold:
                        run.font.bold = True
                    run.font.color.rgb = color
                    n += 1
    return n


def replace_placeholders(prs, mapping: dict[str, str]) -> dict[str, int]:
    """Substitui placeholders {{X}} em qualquer parte do texto dos runs.

    Diferente de replace_run_text (que exige match exato), esta faz
    substituição de substring — apropriada para placeholders dentro de
    strings maiores (ex.: "Recife, {{DATA}}").

    mapping: {'{{TITULO}}': 'Status do Projeto', '{{AUTOR}}': 'Raniel Silva', ...}
    Retorna {placeholder: count_aplicado}.
    """
    counts: dict[str, int] = {p: 0 for p in mapping}
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for placeholder, value in mapping.items():
                        if placeholder in run.text:
                            n = run.text.count(placeholder)
                            run.text = run.text.replace(placeholder, value)
                            counts[placeholder] += n
    return counts


# Identidade visual UPE/POLI (referência exposta para uso externo)
UPE_RED = RGBColor(0xE0, 0x26, 0x1E)
UPE_NAVY = RGBColor(0x1F, 0x2A, 0x44)
UPE_GREEN = RGBColor(0x2E, 0x7D, 0x5B)
DEFAULT_FONT = "Arial"
