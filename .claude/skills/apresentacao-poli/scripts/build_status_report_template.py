#!/usr/bin/env python3
"""Constrói o template genérico status-report.pptx a partir do template antigo.

Estratégia:
1. Carrega o template existente (que tem identidade visual UPE correta, mas
   conteúdo hardcoded de outra apresentação).
2. Para cada slide, mantém os shapes de identidade (logo, footer, rótulo de
   seção, número de página) e substitui o resto por placeholders limpos.
3. Slides 3-6 (conteúdo): substituídos por layout simples "rótulo + título +
   5 bullets".
4. Slide 1 (capa): substitui textos hardcoded por placeholders.
5. Slide 7 (obrigado): substitui textos hardcoded por placeholders.

Uso:
    python3 build_status_report_template.py <template_antigo> <saida>
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

# Cores UPE/POLI extraídas do template original
UPE_RED = RGBColor(0xE0, 0x26, 0x1E)
NAVY = RGBColor(0x1F, 0x3A, 0x6B)
GRAY_DARK = RGBColor(0x4A, 0x4A, 0x5A)
GRAY_LIGHT = RGBColor(0x6B, 0x6B, 0x7A)


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def set_text_keeping_first_format(shape, new_text: str) -> None:
    """Substitui o texto de um shape (todas as runs/paragraphs) por new_text,
    preservando a formatação do primeiro run encontrado.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    if not first_p.runs:
        # Sem run, apenas seta o texto
        tf.text = new_text
        return
    first_run = first_p.runs[0]

    # Capturar formatação do primeiro run
    font_name = first_run.font.name
    font_size = first_run.font.size
    font_bold = first_run.font.bold
    try:
        font_color = first_run.font.color.rgb
    except (AttributeError, TypeError):
        font_color = None

    # Encontrar txBody (parent dos parágrafos)
    tx_body = first_p._p.getparent()
    # Encontrar APENAS os elementos <a:p> (não bodyPr, lstStyle, etc.)
    all_p = tx_body.findall(f"{{{_NS_A}}}p")
    # Deletar todos os parágrafos exceto o primeiro
    for p_elem in all_p[1:]:
        tx_body.remove(p_elem)
    # Limpar runs do primeiro parágrafo
    for run in list(first_p.runs):
        first_p._p.remove(run._r)

    # Adicionar novo run
    new_run = first_p.add_run()
    new_run.text = new_text
    if font_name:
        new_run.font.name = font_name
    if font_size:
        new_run.font.size = font_size
    if font_bold is not None:
        new_run.font.bold = font_bold
    if font_color:
        new_run.font.color.rgb = font_color


def replace_run_text_exact(slide, old: str, new: str) -> bool:
    """Substitui o texto exato de um run em um slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old:
                    run.text = new
                    return True
    return False


def find_shape_by_text(slide, text: str):
    """Retorna o shape cujo primeiro run contém o texto."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if text in run.text:
                    return shape
    return None


def get_shape_first_text(shape) -> str:
    """Retorna o primeiro texto não vazio de um shape."""
    if not shape.has_text_frame:
        return ""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                return run.text
    return ""


def fix_slide_1_capa(slide) -> None:
    """Capa: substitui hardcoded title + nome autor + label orientador + nome orientador
    por placeholders limpos. Preserva: logo UPE, retângulo decorativo, todos os outros
    placeholders.
    """
    # Título principal hardcoded ("Agentes Conversacionais...") → {{TITULO_PRINCIPAL}}
    title_shape = find_shape_by_text(slide, "Agentes Conversacionais")
    if title_shape:
        set_text_keeping_first_format(title_shape, "{{TITULO_PRINCIPAL}}")

    # Autor: "Raniel Silva" → {{AUTOR}}
    replace_run_text_exact(slide, "Raniel Silva", "{{AUTOR}}")
    # Rótulo "Orientador" → {{ROTULO_ORIENTADOR}}
    replace_run_text_exact(slide, "Orientador", "{{ROTULO_ORIENTADOR}}")
    # Nome orientador → {{ORIENTADOR}}
    replace_run_text_exact(slide, "Prof. Dr. Alexandre Maciel", "{{ORIENTADOR}}")
    # Local + data → {{LOCAL_DATA}}
    replace_run_text_exact(slide, "Recife, {{DATA}}", "{{LOCAL_DATA}}")


def fix_slide_2_agenda(slide) -> None:
    """Slide 2 (agenda): ajusta página '2 / 18' → '2 / 7' e remove a linha
    'Duração estimada: {{DURACAO}}' (informação técnica desnecessária na agenda)."""
    replace_run_text_exact(slide, "2 / 18", "2 / 7")
    # Remover textbox com "Duração estimada"
    shape = find_shape_by_text(slide, "Duração estimada")
    if shape:
        shape._element.getparent().remove(shape._element)


def rebuild_content_slide(slide, slide_number: int, total_slides: int = 7) -> None:
    """Reconstrói slide de conteúdo (3-6) com layout simples:

    - Mantém: logo (Picture), rótulo de seção (topo), footer, número de página
    - Remove: todos os outros shapes
    - Adiciona: textbox de título ({{SLIDE_TITULO}}) e textbox de bullets
    """
    # Identificar shapes a PRESERVAR pelo XML element (estável entre iterações)
    elements_to_keep: set = set()
    rotulo_shape_xml = None
    pagina_text_old = None

    for shape in slide.shapes:
        # Logo: Picture com left > ~10 inches (canto superior direito)
        if shape.shape_type == 13:
            try:
                if shape.left and Emu(shape.left).inches > 10:
                    elements_to_keep.add(shape._element)
                    continue
            except Exception:
                pass

        if not shape.has_text_frame:
            continue

        first_text = get_shape_first_text(shape)
        if not first_text:
            continue

        # Rótulo de seção: TextBox pequena no topo (T < 0.7) sem placeholder
        try:
            top_in = Emu(shape.top).inches if shape.top else 99
        except Exception:
            top_in = 99
        if top_in < 0.7 and "{{" not in first_text:
            rotulo_shape_xml = shape._element
            elements_to_keep.add(shape._element)
            continue

        # Footer
        if first_text == "{{FOOTER}}":
            elements_to_keep.add(shape._element)
            continue

        # Número de página (formato "N / 18")
        if " / 18" in first_text:
            elements_to_keep.add(shape._element)
            pagina_text_old = first_text
            continue

    # Remover shapes que não estão na lista de preservados
    to_remove = []
    for shape in slide.shapes:
        if shape._element not in elements_to_keep:
            to_remove.append(shape._element)

    for elem in to_remove:
        elem.getparent().remove(elem)

    # Substituir rótulo por {{SLIDE_ROTULO}}
    if rotulo_shape_xml is not None:
        # Re-localizar o shape após remoção
        for shape in slide.shapes:
            if shape._element == rotulo_shape_xml:
                set_text_keeping_first_format(shape, "{{SLIDE_ROTULO}}")
                break

    # Atualizar número de página
    if pagina_text_old:
        replace_run_text_exact(slide, pagina_text_old, f"{slide_number} / {total_slides}")

    # Adicionar TÍTULO do slide ({{SLIDE_TITULO}})
    title_box = slide.shapes.add_textbox(
        left=Inches(0.5), top=Inches(0.9), width=Inches(11.0), height=Inches(0.8)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "{{SLIDE_TITULO}}"
    title_run.font.name = "Arial"
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = NAVY

    # Adicionar BULLETS ({{BULLET_1}}..{{BULLET_3}}) — 3 bullets por slide,
    # fonte maior (22pt) e espaçamento generoso para visual mais limpo
    bullets_box = slide.shapes.add_textbox(
        left=Inches(0.7), top=Inches(2.3), width=Inches(12.0), height=Inches(4.5)
    )
    bullets_tf = bullets_box.text_frame
    bullets_tf.word_wrap = True

    for i in range(1, 4):  # apenas 3 bullets
        if i == 1:
            p = bullets_tf.paragraphs[0]
        else:
            p = bullets_tf.add_paragraph()
        run = p.add_run()
        run.text = f"•   {{{{BULLET_{i}}}}}"
        run.font.name = "Arial"
        run.font.size = Pt(22)
        run.font.color.rgb = GRAY_DARK
        p.space_after = Pt(24)  # mais respiro entre bullets


def fix_slide_7_obrigado(slide) -> None:
    """Slide final: substitui hardcoded por placeholders genéricos."""
    # "Discussão e perguntas" → {{MENSAGEM_FINAL}}
    replace_run_text_exact(slide, "Discussão e perguntas", "{{MENSAGEM_FINAL}}")
    # Remover textos hardcoded da documentação
    for text in (
        "Documentação completa disponível no repositório do projeto",
        "Resultados da RSL  ·  Framework FMD-Agent  ·  Worklog",
    ):
        shape = find_shape_by_text(slide, text)
        if shape:
            shape._element.getparent().remove(shape._element)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("template_antigo", type=Path)
    ap.add_argument("saida", type=Path)
    args = ap.parse_args()

    prs = Presentation(args.template_antigo)

    print(f"Template carregado: {len(prs.slides)} slides")

    print("Editando slide 1 (capa)...")
    fix_slide_1_capa(prs.slides[0])

    print("Editando slide 2 (agenda)...")
    fix_slide_2_agenda(prs.slides[1])

    for i in range(2, 6):
        print(f"Reconstruindo slide {i+1} (conteúdo)...")
        rebuild_content_slide(prs.slides[i], slide_number=i+1, total_slides=7)

    print("Editando slide 7 (obrigado)...")
    fix_slide_7_obrigado(prs.slides[6])

    args.saida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.saida)
    print(f"\n✓ Salvo: {args.saida}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
