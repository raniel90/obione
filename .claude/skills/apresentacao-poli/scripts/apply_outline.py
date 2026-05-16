#!/usr/bin/env python3
"""Aplica um outline em Markdown sobre um template .pptx, gerando novo .pptx.

Uso:
    python3 apply_outline.py <outline.md> <template.pptx> <saida.pptx>

Fluxo:
    1. Parseia o outline (frontmatter YAML-like + seções markdown)
    2. Substitui placeholders globais no template ({{TITULO}}, {{AUTOR}}, etc.)
    3. Para cada slide do outline, substitui o slide correspondente do template
       por título e bullets (match por id ou ordem)
    4. Remove slides não usados (declarados como 'remover' ou em excesso)
    5. Atualiza footers "X / N"
    6. Salva no caminho de saída

Formato do outline (ver references/outline-format.md):

    ---
    titulo: Status do Projeto Orion
    autor: Raniel Silva
    programa: Doutorado PPGEC
    instituicao: UPE/POLI
    data: 2026-05-15
    template: status-report
    ---

    # slide:capa
    título: Status Report — Maio/2026

    # slide:agenda
    1: Contexto
    2: Andamento
    3: Próximos passos
    4: Riscos

    # slide:conteudo id=1
    titulo: Contexto do projeto
    bullets:
      - Objetivo: ...
      - Stakeholders: ...
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation

# Permite importar pptx_helpers do mesmo diretório
sys.path.insert(0, str(Path(__file__).parent))
from pptx_helpers import (  # noqa: E402
    delete_slide_by_index,
    replace_placeholders,
    replace_run_text,
    update_footers,
)


def _replace_text_in_slide(slide, old: str, new: str) -> int:
    """Substitui substring `old` por `new` em qualquer run de qualquer shape do slide.

    Diferente de `replace_run_text` (match exato), suporta runs que contenham
    o placeholder dentro de outra string — útil quando o template tem prefixos
    como "•   {{BULLET_1}}".
    """
    n = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    n += run.text.count(old)
                    run.text = run.text.replace(old, new)
    return n


def parse_outline(md_text: str) -> tuple[dict, list[dict]]:
    """Parser simples: frontmatter (entre ---) + seções (# slide:tipo).

    Retorna (metadata, [slide_specs]).
    """
    lines = md_text.splitlines()
    meta: dict = {}
    slides: list[dict] = []
    i = 0

    # Frontmatter
    if lines and lines[0].strip() == "---":
        i = 1
        while i < len(lines) and lines[i].strip() != "---":
            line = lines[i].strip()
            if ":" in line:
                k, _, v = line.partition(":")
                meta[k.strip()] = v.strip()
            i += 1
        i += 1  # skip closing ---

    # Seções de slide
    current: dict | None = None
    while i < len(lines):
        line = lines[i]
        m = re.match(r"^#\s+slide:(\w+)(?:\s+id=(\S+))?", line.strip())
        if m:
            if current:
                slides.append(current)
            current = {"tipo": m.group(1), "id": m.group(2), "campos": {}, "bullets": []}
        elif current is not None:
            stripped = line.strip()
            if stripped.startswith("- "):
                current["bullets"].append(stripped[2:].strip())
            elif ":" in stripped and not stripped.startswith("#"):
                k, _, v = stripped.partition(":")
                current["campos"][k.strip()] = v.strip()
        i += 1
    if current:
        slides.append(current)

    return meta, slides


def build_footer(meta: dict) -> str:
    """Monta o footer no padrão 'Nome · Programa · Instituição · Mês Ano'."""
    parts = [
        meta.get("autor", "{{AUTOR}}"),
        meta.get("programa", "{{PROGRAMA}}"),
        meta.get("instituicao", "UPE/POLI"),
        _format_month_year(meta.get("data", "{{DATA}}")),
    ]
    return " · ".join(parts)


def _format_month_year(date_str: str) -> str:
    """Converte 'YYYY-MM-DD' em 'Mês YYYY' (PT-BR). Se falhar, retorna como veio."""
    months_pt = {
        1: "Janeiro", 2: "Fevereiro", 3: "Março", 4: "Abril",
        5: "Maio", 6: "Junho", 7: "Julho", 8: "Agosto",
        9: "Setembro", 10: "Outubro", 11: "Novembro", 12: "Dezembro",
    }
    m = re.match(r"(\d{4})-(\d{2})", date_str)
    if not m:
        return date_str
    year, month = int(m.group(1)), int(m.group(2))
    return f"{months_pt.get(month, month)} {year}"


def apply(meta: dict, slide_specs: list[dict], template_path: Path, out_path: Path) -> None:
    """Aplica o outline ao template e salva."""
    prs = Presentation(template_path)

    # 1. Placeholders globais (defaults razoáveis quando outline não fornece)
    placeholder_map = {
        "{{TITULO}}": meta.get("titulo", ""),
        "{{TITULO_PRINCIPAL}}": meta.get("titulo_principal", meta.get("titulo", "")),
        "{{AUTOR}}": meta.get("autor", ""),
        "{{PROGRAMA}}": meta.get("programa", ""),
        "{{INSTITUICAO}}": meta.get("instituicao", "UPE/POLI"),
        "{{DATA}}": _format_month_year(meta.get("data", "")),
        "{{FOOTER}}": build_footer(meta),
        # Defaults para placeholders herdados dos templates
        "{{CONTEXTO}}": meta.get("contexto", meta.get("titulo", "")).upper(),
        "{{ROTULO_AUTOR}}": meta.get("rotulo_autor", "Autor"),
        "{{ROTULO_ORIENTADOR}}": meta.get("rotulo_orientador", "Orientador"),
        "{{ORIENTADOR}}": meta.get("orientador", ""),
        "{{PROGRAMA_LONGO}}": meta.get(
            "programa_longo",
            meta.get("programa", "Programa de Pós-Graduação"),
        ),
        "{{INSTITUICAO_LONGA}}": meta.get(
            "instituicao_longa",
            "Universidade de Pernambuco · Escola Politécnica de Pernambuco",
        ),
        "{{LOCAL_DATA}}": meta.get(
            "local_data",
            f"Recife, {_format_month_year(meta.get('data', ''))}",
        ),
        "{{ROTULO_AGENDA}}": meta.get("rotulo_agenda", "PLANO DA REUNIÃO"),
        "{{DURACAO}}": meta.get("duracao", "~15 minutos + discussão"),
        "{{MENSAGEM_FINAL}}": meta.get("mensagem_final", "Discussão e perguntas"),
    }
    counts = replace_placeholders(prs, placeholder_map)
    print("Placeholders globais substituídos:")
    for k, n in counts.items():
        print(f"  {k}: {n}×")

    # 2. Aplicar cada slide do outline ao slide do template de mesma posição
    #    Estratégia simples: outline[i] mapeia para prs.slides[i] do template.
    n_template = len(prs.slides)
    n_outline = len(slide_specs)
    print(f"\nSlides no template: {n_template}, no outline: {n_outline}")

    for i, spec in enumerate(slide_specs):
        if i >= n_template:
            print(f"  ⚠️  Slide {i+1} do outline sem correspondente no template — ignorado")
            continue
        slide = prs.slides[i]
        # Substituir título se especificado (substring match para suportar prefixos no run)
        titulo = spec["campos"].get("titulo") or spec["campos"].get("título")
        if titulo:
            for marker in ("{{SLIDE_TITULO}}", "{{TITULO_SLIDE}}", "{{TITLE}}"):
                if _replace_text_in_slide(slide, marker, titulo):
                    break
        # Substituir bullets — substring match (runs podem ter prefixos tipo "•   ")
        for j, bullet in enumerate(spec["bullets"], start=1):
            marker = f"{{{{BULLET_{j}}}}}"
            _replace_text_in_slide(slide, marker, bullet)
        # Campos extras (rótulo, agenda items, etc.) — substring match para flexibilidade
        for k, v in spec["campos"].items():
            _replace_text_in_slide(slide, f"{{{{{k.upper()}}}}}", v)

    # 3. Remover slides excedentes do template
    if n_outline < n_template:
        for idx in range(n_template - 1, n_outline - 1, -1):
            delete_slide_by_index(prs, idx)
        print(f"  Removidos {n_template - n_outline} slides em excesso do template")

    # 4. Atualizar footers
    n_foot = update_footers(prs)
    print(f"\nFooters atualizados: {n_foot}")

    # 5. Salvar
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"\n✓ Salvo: {out_path}")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("outline", type=Path, help="caminho do outline.md")
    ap.add_argument("template", type=Path, help="caminho do template.pptx")
    ap.add_argument("saida", type=Path, help="caminho do .pptx de saída")
    args = ap.parse_args()

    if not args.outline.exists():
        print(f"outline não encontrado: {args.outline}", file=sys.stderr)
        return 2
    if not args.template.exists():
        print(f"template não encontrado: {args.template}", file=sys.stderr)
        return 2

    md = args.outline.read_text(encoding="utf-8")
    meta, slide_specs = parse_outline(md)
    print(f"Meta: {meta}\nSlides especificados: {len(slide_specs)}\n")
    apply(meta, slide_specs, args.template, args.saida)
    return 0


if __name__ == "__main__":
    sys.exit(main())
